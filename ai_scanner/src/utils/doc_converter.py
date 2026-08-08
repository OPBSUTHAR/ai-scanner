import csv, os
from pathlib import Path

OFFICE_EXT = {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv", ".tsv", ".pptx", ".ppt", ".txt", ".md", ".log", ".rtf"}


def is_office_file(name: str) -> bool:
    return Path(name).suffix.lower() in OFFICE_EXT


def convert_to_pdf(src: str) -> str:
    src_p = Path(src)
    ext = src_p.suffix.lower()
    if ext == ".pdf":
        return str(src_p)
    out = src_p.with_suffix(".pdf")
    try:
        if ext in (".docx",):
            _docx_to_pdf(src_p, out)
        elif ext in (".xlsx", ".xls"):
            _xlsx_to_pdf(src_p, out)
        elif ext in (".csv", ".tsv"):
            _csv_to_pdf(src_p, out)
        elif ext in (".pptx", ".ppt"):
            _pptx_to_pdf(src_p, out)
        elif ext in (".txt", ".md", ".log", ".rtf"):
            _text_to_pdf(src_p, out)
        else:
            return ""
        return str(out) if out.exists() else ""
    except Exception:
        return ""


def _build_doc(title, elements):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(title), pagesize=letter,
                            topMargin=0.6 * inch, bottomMargin=0.6 * inch,
                            leftMargin=0.6 * inch, rightMargin=0.6 * inch)
    doc.build(elements)
    return title


def _bank_lines(text, styles):
    from reportlab.platypus import Paragraph, Spacer
    from reportlab.lib.units import inch
    els = []
    for line in text.splitlines() or [""]:
        if not line.strip():
            els.append(Spacer(1, 0.08 * inch))
        else:
            els.append(Paragraph(line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), styles["BodyText"]))
    return els


def _docx_to_pdf(src: Path, out: Path):
    from docx import Document
    from reportlab.platypus import Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    styles = getSampleStyleSheet()
    doc = Document(str(src))
    els = []
    for p in doc.paragraphs:
        txt = (p.text or "").strip()
        if not txt:
            continue
        els.append(Paragraph(txt.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), styles["BodyText"]))
    for table in doc.tables[:50]:
        rows = []
        for row in table.rows:
            rows.append([(c.text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;") for c in row.cells])
        if rows:
            t = Table(rows, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e5decd")),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]))
            els.append(Spacer(1, 0.1 * inch))
            els.append(t)
    _text_doc = _build_doc(out, els)


def _xlsx_to_pdf(src: Path, out: Path):
    from openpyxl import load_workbook
    from reportlab.platypus import Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    styles = getSampleStyleSheet()
    wb = load_workbook(str(src), read_only=True, data_only=True)
    els = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append([("" if v is None else str(v)).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")[:300] for v in row])
        if rows:
            t = Table(rows, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#C7A96E")),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                ("FONTSIZE", (0, 0), (-1, -1), 4),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 2),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
            ]))
            els.append(Paragraph(f"<b>{ws.title}</b>", styles["Heading6"]))
            els.append(Spacer(1, 0.1 * inch))
            els.append(t)
            els.append(PageBreak())
    _build_doc(out, els)


def _csv_to_pdf(src: Path, out: Path):
    from reportlab.platypus import Table, TableStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    import csv
    with open(src, "r", encoding="utf-8", errors="replace") as f:
        reader = list(csv.reader(f))
    table = [[(c or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")[:300] for c in row] for row in reader]
    t = Table(table, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#C7A96E")),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
        ("FONTSIZE", (0, 0), (-1, -1), 4),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    _build_doc(out, [t])


def _pptx_to_pdf(src: Path, out: Path):
    from pptx import Presentation
    from reportlab.platypus import Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    styles = getSampleStyleSheet()
    prs = Presentation(str(src))
    els = []
    for i, slide in enumerate(prs.slides, 1):
        els.append(Paragraph(f"<b>Slide {i}</b>", styles["Heading6"]))
        els.append(Spacer(1, 0.08 * inch))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = "".join(r.text or "" for r in p.runs) or (p.text or "")
                    txt = txt.strip()
                    if not txt:
                        continue
                    els.append(Paragraph(txt.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), styles["BodyText"]))
        els.append(PageBreak())
    _build_doc(out, els)


def _text_to_pdf(src: Path, out: Path):
    from reportlab.lib.styles import getSampleStyleSheet
    text = src.read_text(encoding="utf-8", errors="replace")
    _build_doc(out, _bank_lines(text, getSampleStyleSheet()))


def merge_pdfs(pdf_paths: list, out_path: str) -> str:
    from pypdf import PdfWriter
    writer = PdfWriter()
    for p in pdf_paths:
        try:
            writer.append(str(p))
        except Exception:
            continue
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    with open(out, "wb") as f:
        writer.write(f)
    return str(out)